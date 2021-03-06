#!/bin/env python
# this module should be called lalaloopy
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import sys

def table_size(data, size):
    for i in xrange(0, len(data), size):
        yield data[i:i + size]
        
def table_create(prs,position,title_text,table_header,table_data,cell_widths):
  try: 
    cols = position['cols']
    rows = position['rows']
    top = Inches(position['top'])
    left = Inches(position['left'])
    width = Inches(position['width'])
    height = Inches(position['height'])

    for data in table_size(table_data,rows):
      title_only_slide_layout = prs.slide_layouts[5]
      slide = prs.slides.add_slide(title_only_slide_layout)
      shapes = slide.shapes
      shapes.title.text = title_text
      # +1 for header
      table = shapes.add_table(rows+1, cols, left, top, width, height).table
    
      for i in xrange(0, len(cell_widths)):
        table.columns[i].width = Inches(cell_widths[i])
      
      data.insert(0,table_header)
      for y in xrange(0, len(data)):
        for x in xrange(0, cols):
          table.cell(y, x).text = data[y][x]
          table.cell(y, x).text_frame.paragraphs[0].font.size = Pt(12)
  except:
    pass
    #print sys.exc_info()[0]
    #raise

def main():
  module = AnsibleModule(
    argument_spec=dict(
      filename = dict(required=True),
      title = dict(required=True),
      position = dict(required=False, type='dict', default={'cols': 5, 'rows': 9, 'left': 0.6 , 'top': 2.3, 'width': 10.5, 'height': 0.8 }),
      table_data = dict(required=True, type='list'),
      table_header = dict(required=False, type='list', default=['Name','Project','Role','Ansible Exp','Email']),
      cell_widths = dict(required=False, type='list', default=[2.0,2.0,3.0,2.0,3.0])
      ),
      supports_check_mode=False
  )
  try:
    table_data = module.params['table_data']
    title_text = module.params['title']
    filename = module.params['filename']
    position = module.params['position']
    table_header = module.params['table_header']
    cell_widths = module.params['cell_widths']
    rows = position['rows']

    prs = Presentation(filename)   
    table_create(prs,position,title_text,table_header,table_data,cell_widths)
    prs.save(filename)
    module.exit_json(changed=True)
  except:
    pass
    #print sys.exc_info()[0]
    #raise

from ansible.module_utils.basic import *
main()
