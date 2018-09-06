#!/bin/env python
from pptx import Presentation
from pptx.util import Inches
import sys


def table_size(data, size):
    for i in range(0, len(data), size):
        yield data[i:i + size]
        
def table_create(prs,position,title_text,table_header,table_data):
  try: 
    cols = position['cols']
    rows = position['rows']
    left = top = Inches(position['left'])
    width = Inches(position['width'])
    height = Inches(position['height'])

    for data in table_size(table_data,rows):
      title_only_slide_layout = prs.slide_layouts[5]
      slide = prs.slides.add_slide(title_only_slide_layout)
      shapes = slide.shapes
      shapes.title.text = title_text
      table = shapes.add_table(rows, cols, left, top, width, height).table
    
      for i in xrange(0, len(table_header)):
        table.columns[i].width = Inches(table_header[i])
      
      for y in xrange(0, len(data)):
        for x in xrange(0, cols):
          table.cell(y, x).text = data[y][x]
  except:
    print sys.exc_info()[0]
    raise

def main():
  module = AnsibleModule(
    argument_spec=dict(
      filename = dict(required=True),
      title = dict(required=True),
      position = dict(required=False, type='dict', default={'cols': 5, 'rows': 20, 'left': 1.0 , 'top': 2.0, 'width': 8.0, 'height': 0.8 }),
      table_data = dict(required=True, type='list'),
      table_header = dict(required=False, type='list', default=[2.0,1.0,2.0,1.0,2.0])
      #table_header = dict(required=False, type='dict', default={'Name': 2.0, 'Project': 1.0, 'Role': 2.0 , 'Ansible': 1.0, 'Email': 2.0 })
      ),
      supports_check_mode=False
  )
  try:
    table_data = module.params['table_data']
    title_text = module.params['title']
    filename = module.params['filename']
    position = module.params['position']
    table_header = module.params['table_header']
    rows = position['rows']

    prs = Presentation(filename)   
    #for i in xrange(0, len(table_data), rows):
    #  table_create(prs,position,title_text,table_header,table_data[i:i + rows])
    table_create(prs,position,title_text,table_header,table_data)
    prs.save(filename)
    module.exit_json(changed=True)
  except:
    print sys.exc_info()[0]
    raise

from ansible.module_utils.basic import *
main()
