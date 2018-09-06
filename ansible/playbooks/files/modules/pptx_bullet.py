#!/bin/env python
from pptx import Presentation
import sys
from pptx import Presentation

def bullet_size(data, size):
    for i in xrange(0, len(data), size):
        yield data[i:i + size]

def main():
  module = AnsibleModule(
    argument_spec=dict(
      filename = dict(required=True),
      title = dict(required=True),
      count = dict(required=False, type='int', default=5),
      bullets = dict(required=True, type='list')
      ),
      supports_check_mode=False
  )
  try:
    filename = module.params['filename']
    prs = Presentation(filename)
    
    bullets = module.params['bullets']
    for data in bullet_size(bullets,count):
      bullet_slide_layout = prs.slide_layouts[1]
      slide = prs.slides.add_slide(bullet_slide_layout)
      shapes = slide.shapes
      body_shape = shapes.placeholders[1]

      title_shape = shapes.title
      title_shape.text = module.params['title']
      tf = body_shape.text_frame
      tf.text = data[0]
    
      for x in xrange(1, len(data)):
        if data[x]: 
          p = tf.add_paragraph()
          p.text = data[x]
          p.level = 0

    prs.save(filename)
    module.exit_json(changed=True)
  except:
    #pass
    print sys.exc_info()[0]
    raise
    #module.fail_json(msg=sys.exc_info()[0])

from ansible.module_utils.basic import *
main()
