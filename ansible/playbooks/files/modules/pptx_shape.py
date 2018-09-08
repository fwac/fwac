#!/bin/env python
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

def bullet_size(data, size):
    for i in xrange(0, len(data), size):
        yield data[i:i + size]

def main():
  module = AnsibleModule(
    argument_spec=dict(
      filename = dict(required=True),
      title = dict(required=True),
      count = dict(required=False, type='int', default=5),
      bullets = dict(required=True, type='list')
      ),
      supports_check_mode=False
  )
  try:
    filename = module.params['filename']
    count = module.params['count']
    prs = Presentation(filename)
    bullets = module.params['bullets']
    


    for shape in bullet_size(shapes,count):
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = 'Adding an AutoShape'

        left = Inches(0.93)  # 0.93" centers this overall set of shapes
        top = Inches(3.0)
        width = Inches(1.75)
        height = Inches(1.0)

        shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
        shape.text = 'Step 1'

        left = left + width - Inches(0.4)
        width = Inches(2.0)  # chevrons need more width for visual balance

        for n in range(2, 6):
            shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
            shape.text = 'Step %d' % n
            left = left + width - Inches(0.4)


        
    for data in bullet_size(bullets,count):
      bullet_slide_layout = prs.slide_layouts[1]
      slide = prs.slides.add_slide(bullet_slide_layout)
      shapes = slide.shapes
      body_shape = shapes.placeholders[1]

      title_shape = shapes.title
      title_shape.text = module.params['title']
      tf = body_shape.text_frame
      tf.text = data[0]
    
      for x in xrange(1, len(data)):
        if data[x]: 
          p = tf.add_paragraph()
          p.text = data[x]
          p.level = 0

    prs.save(filename)
    module.exit_json(changed=True)
  except:
    pass
    #print sys.exc_info()[0]
    #raise
    #module.fail_json(msg=sys.exc_info()[0])

from ansible.module_utils.basic import *
main()
