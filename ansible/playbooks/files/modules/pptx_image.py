#!/usr/bin/python
from pptx import Presentation
from pptx.util import Inches
import sys


def main():
  module = AnsibleModule(
    argument_spec=dict(
      filename = dict(required=True),
      title = dict(required=True),
      top_inches = dict(required=False, type='float', default='.6'),
      left_inches = dict(required=False, type='float', default='1.0'),      
      image = dict(required=True),
      bg_image = dict(required=False)
      ),
      supports_check_mode=False
  )
  try:
    filename = module.params['filename']
    prs = Presentation(filename)
    
    top_inches = module.params['top_inches']
    left_inches = module.params['left_inches']
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    top = Inches(top_inches)
    left = Inches(left_inches)
    if bg_image:
      img_path = module.params['bg_image']
      pic = slide.shapes.add_picture(img_path, left, top)   
    img_path = module.params['image']
    pic = slide.shapes.add_picture(img_path, left, top)    
    prs.save(filename)
    module.exit_json(changed=True)
  except:
    #pass
    print sys.exc_info()[0]
    raise

from ansible.module_utils.basic import *
main()
