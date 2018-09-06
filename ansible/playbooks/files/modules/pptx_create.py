#!/bin/env python
import sys
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def main():
  module = AnsibleModule(
    argument_spec=dict(
      filename = dict(required=True),
      template = dict(required=False),
      title = dict(default=True),
      subtitle = dict(default=True),
      ),
      supports_check_mode=False
  )
  try:
    template = module.params['template']
    if template:
      prs = Presentation(template)
    else:
      prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = module.params['title']
    subtitle.text = module.params['subtitle']
    prs.save(module.params['filename'])
    module.exit_json(changed=True)
  except:
    module.fail_json(msg=sys.exc_info()[0])

from ansible.module_utils.basic import *
main()
