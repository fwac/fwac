#!/usr/bin/python
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.util import Inches
import sys

class FWChart:

  def __init__(self,filename):
    self.prs = Presentation(filename)

  def bar_chart(self,categories,series_name,series_values,title):
    try:
      count_items = []
      for cat in categories:
        count_items.append(series_values.count(cat))
        
      slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
      slide_title = slide.shapes.title
      slide_title.text = title
      chart_data = ChartData()
      chart_data.categories = categories      
      chart_data.add_series(series_name,count_items)
      x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
      slide.shapes.add_chart(
          XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
      )
    except:
      pass

  def pie_chart(self,categories,series_name,series_values,build,title):
    try:
      count_items = []
      if build:
        for cat in categories:
          count_items.append(series_values.count(float(cat)))
      else:
        count_items = series_values
      # convert to percentages  
      total_count = sum(count_items)  
      percent_items = (x/total_count for x in count_items)
      
      slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
      slide_title = slide.shapes.title
      slide_title.text = title
      chart_data = ChartData()
      chart_data.categories = categories
      chart_data.add_series(series_name, percent_items)
      x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5) 
      chart = slide.shapes.add_chart(
          XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
      ).chart

      chart.has_legend = True
      chart.legend.position = XL_LEGEND_POSITION.BOTTOM
      chart.legend.include_in_layout = False

      chart.plots[0].has_data_labels = True
      data_labels = chart.plots[0].data_labels
      data_labels.number_format = '0%'
      data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    except:
      print sys.exc_info()[0]
      raise

  def save(self,filename):
    try:
      self.prs.save(filename)
    except:
      pass
 
def main():
  module = AnsibleModule(
    argument_spec=dict(
      filename = dict(required=True),
      categories = dict(required=True, type='list'),
      series_name = dict(required=True),
      series_build = dict(required=False, type='bool', default=True),
      series_values = dict(required=True, type='list'),
      title = dict(required=True, type='str'),
      chart_type = dict(required=True, type='str', choices=['bar', 'pie'])
      ),
      supports_check_mode=False
  )
  filename = module.params['filename']
  categories = module.params['categories'] #list
  series_name = module.params['series_name'] 
  title = module.params['title'] 
  series_build = module.params['series_build']
  series_values = module.params['series_values'] #list to tuple
  chart_type = module.params['chart_type']
  try:
    chartmp = FWChart(filename)
    if chart_type == 'bar':
      chartmp.bar_chart(categories,series_name,series_values,title)
    elif chart_type == 'pie':
      chartmp.pie_chart(categories,series_name,series_values,series_build,title)
    chartmp.save(filename)
    module.exit_json(changed=True)
  except:
    #print sys.exc_info()[0]
    #raise
    pass
    #module.fail_json(msg="error")

from ansible.module_utils.basic import *
main()
